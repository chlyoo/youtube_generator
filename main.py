from datetime import datetime
import os

#1.1 Searching for images
#1.2 Saving images
#1.3 Duplication Prevention
today = str(datetime.now().date())
#1.Crawling Part
if not os.path.exists(f"./data/{today}"):
    os.makedirs(f"./data/{today}")

#2. Generating PPT
#2.1 Loading Images to each Slide
#2.---––––––––––––––––––––––––––––––––––––––––
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
#디렉토리에서 사진 가져오기 
for item in os.listdir(f"./data/{today}"):
	img_path =os.getcwd()+'/data/'+today+"/"+item
	#print(img_path)
	blank_slide_layout = prs.slide_layouts[6] 
	slide = prs.slides.add_slide(blank_slide_layout)
	left = top = Inches(1)
	pic = slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.75),
							   width=Inches(9), height=Inches(5))





#img_path = '/data/monty-truth.png'
filename = today+"_utub.pptx"

 #오늘 날짜를 가지고 파일이름 생성
prs.save(filename) #PPT저장 

#2.3 Saving PPT as a video
from convert import ppt_to_mp4

# quality:0-100. The level of quality of the slide. The higher the number, the higher the quality.
quality = 60
# resolution:The resolution of the slide. 480,720,1080...
resolution = 720
# frames: The number of frames per second.
frames = 24
# ppt_path:The ppt/pptx/pptm file path.
#filename='test.pptx'

ppt_path = os.path.abspath(filename)
# mp4_path:The mp4 video save path.
if not os.path.exists("/output"):
    os.makedirs("/output")
mp4_path = os.path.abspath(f"output/{today}.mp4")
#mp4_path = os.path.abspath(f"output/{today}.mp4")

# Require Windows system(Media Player was enabled) and Microsoft Office 2010 or higher.
# Converting ppt into video relies on Windows Media Player. So you need to enable Desktop Experience feature.
# More save types please visit: https://docs.microsoft.com/en-us/office/vba/api/powerpoint.ppsaveasfiletype

# quality:0-100. The level of quality of the slide. The higher the number, the higher the quality.
# quality = 60
# # resolution:The resolution of the slide. 480,720,1080...
# resolution = 720
# # frames: The number of frames per second.
# frames = 24

# # ppt_path:The ppt/pptx/pptm file path.
# ppt_path = os.path.abspath('./test.pptx')
# # mp4_path:The mp4 video save path.
# mp4_path = os.path.abspath('./test.mp4')

# ie_temp_dir:The convert cache file path.
# The default path (hidden) is 'C:/Users/username/AppData/Local/Microsoft/Windows/Temporary Internet Files/Content.MSO/ppt'.
# Or 'C:/Users/username/AppData/Local/Microsoft/Windows/INetCache/Content.MSO/ppt'
# You can find the cache folde at IE setting.
# If you don't want clear cache files,assign ie_temp_dir with empty string.
#ie_temp_dir = 'C:/Users/username/AppData/Local/Microsoft/Windows/INetCache/Content.MSO/ppt'
ie_temp_dir = ''

# status:Converting result. 0:failed. -1: timeout. 1:success.
status = 0
# timeout: Seconds that converting time out.
timeout = 4*600000
duration=10
try:
    status = ppt_to_mp4(ppt_path,mp4_path,duration,resolution,frames,quality,timeout)
    # Clear PowerPoint cache after convert completed. When you converted hundreds of files, the cache folder will be huge.
    if ie_temp_dir != '':
        shutil.rmtree(ie_temp_dir, ignore_errors=True)
except Exception as e:
    print('Error! Code: {c}, Message, {m}'.format(c = type(e).__name__, m = str(e)))
    
if status == -1:
    print('Failed:timeout.')
elif status == 1:
    print('Success!')
else:
    if os.path.exists(mp4_path):
        os.remove(mp4_path)
    print('Failed:The ppt may have unknown elements. You can try to convert it manual.')
#3. Youtube Upload
title= '영상제목'
descr = '영상내용'
kwords = '영상 키워드, 컴마로 구분'
categ = 22 #영상 카테고리 
"""
2 - Autos & Vehicles
1 -  Film & Animation
10 - Music
15 - Pets & Animals
17 - Sports
18 - Short Movies
19 - Travel & Events
20 - Gaming
21 - Videoblogging
22 - People & Blogs
23 - Comedy
24 - Entertainment
25 - News & Politics
26 - Howto & Style
27 - Education
28 - Science & Technology
29 - Nonprofits & Activism
30 - Movies
31 - Anime/Animation
32 - Action/Adventure
33 - Classics
34 - Comedy
35 - Documentary
36 - Drama
37 - Family
38 - Foreign
39 - Horror
40 - Sci-Fi/Fantasy
41 - Thriller
42 - Shorts
43 - Shows
44 - Trailers
"""
prvacy = 'public'#영상 공개여부 public, private, unlisted 중 하나
upllll = f'python upload_video.py --file="{mp4_path}"" --title="{title}" --description="{descr}" --keywords="{kwords}" --category="{categ}" --privacyStatus="{prvacy}"'
#3.1 Login to YOUTUBE
#3.2 Upload Generated video to YOUTUBE
#동영상으로 변환
os.system(upllll)